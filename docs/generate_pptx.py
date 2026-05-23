"""
Генератор PPTX-презентации ВКР — Recall (v2 - modern design)
Запуск: source .venv/bin/activate && python docs/generate_pptx.py
Результат: docs/ВКР_Миняев_МИ.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
import copy

BLUE = RGBColor(0x00, 0x3D, 0xA5)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x7A)
ACCENT = RGBColor(0x00, 0x66, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x77)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0xAA)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFC)
CARD_BORDER = RGBColor(0xE0, 0xE5, 0xF0)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

SLIDE_WIDTH = Cm(25.4)
SLIDE_HEIGHT = Cm(19.05)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.5)
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = DARK_BLUE
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = ACCENT
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_accent_dot(slide, x, y):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Cm(0.3), Cm(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Cm(2.2), SLIDE_HEIGHT - Cm(1.3), Cm(1.8), Cm(0.8)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Cm(2), Cm(1.2), Cm(21.5), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    # underline accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.6), Cm(4), Cm(0.12)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_card(slide, left, top, width, height, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(0.5)
    return shape


def add_accent_card(slide, left, top, width, height, accent_color=BLUE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = CARD_BG
    bg.line.color.rgb = CARD_BORDER
    bg.line.width = Pt(0.5)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Cm(0.2), Cm(0.25), height - Cm(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    return bg


def add_text_box(slide, left, top, width, height, text, size=Pt(14),
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        p.space_after = Pt(4)
    return tf


def add_placeholder_box(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xFA)
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xCC)
    shape.line.dash_style = 2
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_metric_card(slide, left, top, width, height, value, desc, accent=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.75)
    # top accent line
    topline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Cm(0.3), top, width - Cm(0.6), Cm(0.15)
    )
    topline.fill.solid()
    topline.fill.fore_color.rgb = accent
    topline.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # =========== СЛАЙД 1: Титульный ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)

    # decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(19), Cm(12), Cm(8), Cm(8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
    circle.line.fill.background()

    add_text_box(slide, Cm(2), Cm(1.8), Cm(21), Cm(2.5),
                 "МИНОБРНАУКИ РОССИИ\n"
                 "ФГБОУ ВО «МИРЭА — Российский технологический университет»\n"
                 "Институт информационных технологий",
                 size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Cm(2), Cm(5.5), Cm(21), Cm(5),
                 "Кроссплатформенное мобильное\nприложение для семантического\n"
                 "анализа медиатеки с использованием\nвстраиваемых LLM",
                 size=Pt(24), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Cm(2), Cm(12.5), Cm(21), Cm(1.2),
                 "студент группы ИКБО-06-22  Миняев М.И.",
                 size=Pt(13), color=BLACK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Cm(2), Cm(14), Cm(21), Cm(1.2),
                 "руководитель: [ФИО руководителя]",
                 size=Pt(11), color=GRAY, alignment=PP_ALIGN.CENTER)

    # =========== СЛАЙД 2: Актуальность ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Актуальность")
    add_slide_number(slide, 2)

    facts = [
        ("4,88 млрд", "пользователей смартфонов (2024)", BLUE),
        ("952+", "фотографий на устройстве в среднем", ACCENT),
        ("70+ TOPS", "мощность мобильных NPU", PURPLE),
        ("0", "офлайн-аналогов на рынке РФ", GREEN),
    ]
    for i, (val, label, color) in enumerate(facts):
        y = Cm(3.3 + i * 3.0)
        add_accent_card(slide, Cm(1.5), y, Cm(11.5), Cm(2.5), color)
        add_text_box(slide, Cm(2.3), y + Cm(0.3), Cm(10), Cm(1.2),
                     val, size=Pt(18), bold=True, color=color)
        add_text_box(slide, Cm(2.3), y + Cm(1.4), Cm(10), Cm(0.8),
                     label, size=Pt(11), color=GRAY)

    # Right side - problem statement
    add_accent_card(slide, Cm(14), Cm(3.3), Cm(10), Cm(3.5), RED_SOFT)
    add_text_box(slide, Cm(15), Cm(3.6), Cm(8.5), Cm(0.8),
                 "Проблема", size=Pt(11), bold=True, color=RED_SOFT)
    add_text_box(slide, Cm(15), Cm(4.4), Cm(8.5), Cm(2),
                 "Пользователь помнит содержимое фото,\nно ищет по дате и имени файла",
                 size=Pt(12), color=BLACK)

    add_accent_card(slide, Cm(14), Cm(7.3), Cm(10), Cm(3.5), ORANGE)
    add_text_box(slide, Cm(15), Cm(7.6), Cm(8.5), Cm(0.8),
                 "Ограничение", size=Pt(11), bold=True, color=ORANGE)
    add_text_box(slide, Cm(15), Cm(8.4), Cm(8.5), Cm(2),
                 "Облачные решения нарушают\nконфиденциальность личных данных",
                 size=Pt(12), color=BLACK)

    add_accent_card(slide, Cm(14), Cm(11.3), Cm(10), Cm(3.5), GREEN)
    add_text_box(slide, Cm(15), Cm(11.6), Cm(8.5), Cm(0.8),
                 "Научная новизна", size=Pt(11), bold=True, color=GREEN)
    add_text_box(slide, Cm(15), Cm(12.4), Cm(8.5), Cm(2),
                 "На рынке РФ нет приложений с локальным\nсемантическим поиском по медиатеке",
                 size=Pt(12), color=BLACK)

    # =========== СЛАЙД 3: Цели и задачи ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Цели и задачи ВКР")
    add_slide_number(slide, 3)

    # Goal card
    add_accent_card(slide, Cm(1.5), Cm(3.3), Cm(11), Cm(10), BLUE)
    add_text_box(slide, Cm(2.3), Cm(3.5), Cm(10), Cm(0.8),
                 "ЦЕЛЬ", size=Pt(11), bold=True, color=BLUE)
    add_text_box(slide, Cm(2.3), Cm(4.5), Cm(9.8), Cm(8.5),
                 "Разработать кроссплатформенное мобильное приложение "
                 "для семантического поиска по локальной медиатеке "
                 "с использованием встроенных LLM, обеспечивающее "
                 "интеллектуальный поиск изображений при полной "
                 "автономности и конфиденциальности данных",
                 size=Pt(13), color=BLACK)

    # Tasks
    add_text_box(slide, Cm(13.5), Cm(3.3), Cm(10), Cm(0.8),
                 "ЗАДАЧИ", size=Pt(11), bold=True, color=BLUE)

    tasks = [
        "Анализ предметной области",
        "Обоснование необходимости автоматизации",
        "Формулировка требований к системе",
        "Разработка архитектуры и информ. модели",
        "Реализация программных компонентов",
        "Тестирование программного решения",
    ]
    for i, task in enumerate(tasks):
        y = Cm(4.5 + i * 1.8)
        # number circle
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Cm(13.5), y, Cm(0.9), Cm(0.9)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, Cm(14.8), y + Cm(0.05), Cm(9), Cm(0.9),
                     task, size=Pt(12), color=BLACK)

    # =========== СЛАЙД 4: Системы-аналоги ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Системы-аналоги")
    add_slide_number(slide, 4)

    headers = ["Система", "Семантический\nпоиск", "Полный\nофлайн", "Кросс-\nплатформ."]
    rows_data = [
        ("Google Photos", "✓", "✗", "✓"),
        ("Apple Photos", "◐", "✗", "✗"),
        ("Simple Gallery Pro", "✗", "✓", "✗"),
        ("Immich", "✓", "✗", "✓"),
        ("Ente Photos", "✗", "✗", "✓"),
        ("Recall (наша)", "✓", "✓", "✓"),
    ]

    table_shape = slide.shapes.add_table(
        len(rows_data) + 1, 4, Cm(1.5), Cm(3.5), Cm(22), Cm(10.5)
    )
    table = table_shape.table
    col_widths = [Cm(7.5), Cm(5), Cm(4.5), Cm(5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            if val == "✓":
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif val == "✗":
                p.font.color.rgb = RED_SOFT
            if row_idx == len(rows_data) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
                p.font.bold = True

    add_text_box(slide, Cm(1.5), Cm(14.8), Cm(22), Cm(1.2),
                 "Ни одно решение не объединяет семантический поиск + офлайн + кроссплатформенность",
                 size=Pt(12), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # =========== СЛАЙД 5: Требования (Use Case) ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Требования к системе")
    add_slide_number(slide, 5)

    add_placeholder_box(slide, Cm(1.5), Cm(3.3), Cm(22), Cm(13),
                        "[Вставить: Диаграмма вариантов использования (Use Case)]\n\n"
                        "Источник: рисунок 2.5 из проектной практики\n\n"
                        "Акторы: Пользователь, Система\n"
                        "Прецеденты: семантический поиск, управление альбомами,\n"
                        "фоновая индексация, мониторинг ресурсов")

    # =========== СЛАЙД 6: Архитектура ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Архитектура системы")
    add_slide_number(slide, 6)

    add_placeholder_box(slide, Cm(1.5), Cm(3.3), Cm(15), Cm(13),
                        "[Вставить: Компонентная диаграмма C4]\n\n"
                        "Источник: рисунок 2.3 из проектной практики\n\n"
                        "UI (Compose) → Feature → Core\n"
                        ":core:ml  :core:vector  :core:database  :core:worker")

    # Principles as modern cards
    principles = [
        ("Clean Architecture", BLUE),
        ("Модульная декомпозиция", ACCENT),
        ("0 серверов", GREEN),
        ("MVI (UDF)", PURPLE),
        ("Offline-first", ORANGE),
    ]
    for i, (text, color) in enumerate(principles):
        y = Cm(3.3 + i * 2.6)
        add_accent_card(slide, Cm(17), y, Cm(7), Cm(2.2), color)
        add_text_box(slide, Cm(17.8), y + Cm(0.5), Cm(6), Cm(1.2),
                     text, size=Pt(12), bold=True, color=color)

    # =========== СЛАЙД 7: Информационная модель ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Информационная модель (схема БД)")
    add_slide_number(slide, 7)

    add_placeholder_box(slide, Cm(1.5), Cm(3.3), Cm(14.5), Cm(13),
                        "[Вставить: ER-диаграмма]\n\n"
                        "Источник: рисунок 2.9 из проектной практики\n\n"
                        "6 сущностей: MediaFile, EmbeddingBatch,\n"
                        "Album, AlbumEntry, ExcludedPath, SearchQuery")

    # Key decisions as cards
    decisions = [
        ("Хранение векторов", "В бинарных батч-файлах,\nне в SQLite", BLUE),
        ("Доступ к вектору", "O(1): offset = k * d * 4 bytes", ACCENT),
        ("СУБД", "SQLite через Room 2.8\n(типобезопасная ORM)", GREEN),
    ]
    for i, (title, desc, color) in enumerate(decisions):
        y = Cm(3.3 + i * 4.2)
        add_accent_card(slide, Cm(17), y, Cm(7), Cm(3.6), color)
        add_text_box(slide, Cm(17.8), y + Cm(0.4), Cm(5.8), Cm(0.8),
                     title, size=Pt(10), bold=True, color=color)
        add_text_box(slide, Cm(17.8), y + Cm(1.4), Cm(5.8), Cm(2),
                     desc, size=Pt(11), color=BLACK)

    # =========== СЛАЙД 8: Математическое обеспечение ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Математическое обеспечение")
    add_slide_number(slide, 8)

    # Section 1: CLIP
    y1 = Cm(3.3)
    add_accent_card(slide, Cm(1.5), y1, Cm(22), Cm(3.8), BLUE)
    add_text_box(slide, Cm(2.3), y1 + Cm(0.3), Cm(20), Cm(0.8),
                 "1. Генерация дескрипторов (CLIP)", size=Pt(13), bold=True, color=BLUE)
    # Formula box
    formula_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2.3), y1 + Cm(1.3), Cm(20.5), Cm(2.0)
    )
    formula_shape.fill.solid()
    formula_shape.fill.fore_color.rgb = WHITE
    formula_shape.line.color.rgb = RGBColor(0xE8, 0xEB, 0xF0)
    tf = formula_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    # Use runs for better formatting
    run = p.add_run()
    run.text = "f"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "img"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    run = p.add_run()
    run.text = " : Image \u2192 \u211d"
    run.font.size = Pt(14)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "d"
    run.font.size = Pt(9)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "          f"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "text"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    run = p.add_run()
    run.text = " : Text \u2192 \u211d"
    run.font.size = Pt(14)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "d"
    run.font.size = Pt(9)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "          d = 384 / 512"
    run.font.size = Pt(13)
    run.font.color.rgb = ACCENT
    run.font.bold = True

    # Section 2: Cosine similarity
    y2 = Cm(7.6)
    add_accent_card(slide, Cm(1.5), y2, Cm(22), Cm(3.8), PURPLE)
    add_text_box(slide, Cm(2.3), y2 + Cm(0.3), Cm(20), Cm(0.8),
                 "2. Поиск ближайшего соседа", size=Pt(13), bold=True, color=PURPLE)
    formula_shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2.3), y2 + Cm(1.3), Cm(20.5), Cm(2.0)
    )
    formula_shape2.fill.solid()
    formula_shape2.fill.fore_color.rgb = WHITE
    formula_shape2.line.color.rgb = RGBColor(0xE8, 0xEB, 0xF0)
    tf = formula_shape2.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "sim(\u0071\u0302, \u0076\u0302"
    run.font.size = Pt(15)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "i"
    run.font.size = Pt(9)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = ") = \u0071\u0302 \u00b7 \u0076\u0302"
    run.font.size = Pt(15)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "i"
    run.font.size = Pt(9)
    run.font.color.rgb = BLACK
    run = p.add_run()
    run.text = "          Results = { i  |  sim \u2265 \u03c4 }"
    run.font.size = Pt(15)
    run.font.color.rgb = BLACK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Косинусное сходство нормированных векторов = скалярное произведение"
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY

    # Section 3: HNSW
    y3 = Cm(11.9)
    add_accent_card(slide, Cm(1.5), y3, Cm(22), Cm(4.5), GREEN)
    add_text_box(slide, Cm(2.3), y3 + Cm(0.3), Cm(20), Cm(0.8),
                 "3. HNSW (Hierarchical Navigable Small World)", size=Pt(13), bold=True, color=GREEN)
    formula_shape3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2.3), y3 + Cm(1.3), Cm(10), Cm(2.7)
    )
    formula_shape3.fill.solid()
    formula_shape3.fill.fore_color.rgb = WHITE
    formula_shape3.line.color.rgb = RGBColor(0xE8, 0xEB, 0xF0)
    tf = formula_shape3.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Линейный перебор  "
    run.font.size = Pt(13)
    run.font.color.rgb = RED_SOFT
    run = p.add_run()
    run.text = "O(n)"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RED_SOFT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = "\u2193"
    run.font.size = Pt(16)
    run.font.color.rgb = GRAY
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run = p3.add_run()
    run.text = "HNSW  "
    run.font.size = Pt(13)
    run.font.color.rgb = GREEN
    run = p3.add_run()
    run.text = "O(log n)"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = GREEN

    # HNSW params
    params_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(13), y3 + Cm(1.3), Cm(10), Cm(2.7)
    )
    params_shape.fill.solid()
    params_shape.fill.fore_color.rgb = WHITE
    params_shape.line.color.rgb = RGBColor(0xE8, 0xEB, 0xF0)
    tf = params_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.5)
    p = tf.paragraphs[0]
    p.text = "Многоуровневый навигируемый граф"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK
    p2 = tf.add_paragraph()
    p2.text = ""
    p2.font.size = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = "M = 16    efConstruction = 200    efSearch = 50"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT

    # =========== СЛАЙД 9: Инструментальные средства ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Инструментальные средства разработки")
    add_slide_number(slide, 9)

    tools = [
        ("Kotlin 2.2", "типобезопасность, корутины", BLUE),
        ("KMP", "единая логика Android + iOS", BLUE),
        ("Jetpack Compose", "декларативный UI, Material 3", ACCENT),
        ("Room 2.8", "типобезопасная ORM для SQLite", ACCENT),
        ("TensorFlow Lite", "инференс CLIP на NPU/GPU", PURPLE),
        ("HNSW (Kotlin)", "векторный поиск O(log n)", PURPLE),
        ("WorkManager", "фоновая индексация + батарея", GREEN),
        ("Hilt", "внедрение зависимостей", GREEN),
        ("MVI", "однонаправленный поток данных", ORANGE),
        ("kotlin.test", "единые тесты для обеих платформ", ORANGE),
        ("Android Studio", "IDE + Profiler (Memory, Energy)", GRAY),
        ("Git", "контроль версий", GRAY),
    ]

    for i, (name, why, color) in enumerate(tools):
        col = i % 2
        row = i // 2
        x = Cm(1.5) + col * Cm(12.2)
        y = Cm(3.3) + row * Cm(2.3)

        add_accent_card(slide, x, y, Cm(11.5), Cm(2.0), color)

        # Tool name and description using text box with runs
        txBox = slide.shapes.add_textbox(x + Cm(0.8), y + Cm(0.4), Cm(10.2), Cm(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = color
        run2 = p.add_run()
        run2.text = f"  —  {why}"
        run2.font.size = Pt(10)
        run2.font.color.rgb = GRAY

    # =========== СЛАЙД 10: Программное решение 1 ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Программное решение")
    add_slide_number(slide, 10)

    screenshots_1 = [
        ("[Экран приветствия]\n\nрис. 4.1 тех. практики", "Onboarding"),
        ("[Главный экран + индексация]\n\nрис. 4.2 тех. практики", "Timeline"),
        ("[Экран альбомов]\n\nfeature/albums", "Альбомы"),
    ]
    for i, (text, caption) in enumerate(screenshots_1):
        x = Cm(1.5) + i * Cm(7.8)
        add_placeholder_box(slide, x, Cm(3.3), Cm(7.2), Cm(11.5), text)
        add_text_box(slide, x, Cm(15.2), Cm(7.2), Cm(0.8),
                     f"Рисунок {i+1} — {caption}",
                     size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)

    # =========== СЛАЙД 11: Программное решение 2 ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Программное решение")
    add_slide_number(slide, 11)

    screenshots_2 = [
        ("[Ввод поискового запроса]\n\nрис. 4.4 тех. практики", "Поиск"),
        ("[Результаты поиска]\n\nрис. 4.5 тех. практики", "Результаты"),
        ("[Настройки исключений]\n\nрис. 4.3 тех. практики", "Исключения"),
    ]
    for i, (text, caption) in enumerate(screenshots_2):
        x = Cm(1.5) + i * Cm(7.8)
        add_placeholder_box(slide, x, Cm(3.3), Cm(7.2), Cm(11.5), text)
        add_text_box(slide, x, Cm(15.2), Cm(7.2), Cm(0.8),
                     f"Рисунок {i+4} — {caption}",
                     size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)

    # =========== СЛАЙД 12: Тестирование — бенчмарки ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Тестирование — бенчмарки производительности")
    add_slide_number(slide, 12)

    # Chart placeholders
    add_placeholder_box(slide, Cm(1.5), Cm(3.3), Cm(11), Cm(5.8),
                        "[График: Search Latency p95]\nHNSW vs LinearScan (1K–50K)\n\n"
                        "Источник: docs/charts-real.html")
    add_placeholder_box(slide, Cm(13), Cm(3.3), Cm(11), Cm(5.8),
                        "[График: Recall@10 vs Index Size]\nM=16, efSearch=50\n\n"
                        "Источник: docs/charts-real.html")

    # Metric cards
    metrics = [
        ("< 2 мс", "p95 @ 50K", BLUE),
        ("\u2265 0.95", "Recall@10", PURPLE),
        ("274 v/s", "throughput", GREEN),
        ("16 мс", "serialize 10K", ORANGE),
    ]
    for i, (val, desc, color) in enumerate(metrics):
        x = Cm(1.5) + i * Cm(5.8)
        add_metric_card(slide, x, Cm(9.8), Cm(5.3), Cm(2.8), val, desc, color)

    add_text_box(slide, Cm(1.5), Cm(13.2), Cm(22), Cm(3),
                 "17 бенчмарк-тестов  |  384d  |  seed=42  |  100 queries\n"
                 "Требование ТЗ: < 50 мс @ 50K  \u2014  выполнено с запасом в 25x",
                 size=Pt(11), color=GRAY, alignment=PP_ALIGN.CENTER)

    # =========== СЛАЙД 13: Тестирование — память ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Тестирование — память и функциональные тесты")
    add_slide_number(slide, 13)

    # Memory section
    add_text_box(slide, Cm(1.5), Cm(3.2), Cm(11), Cm(0.8),
                 "Потребление памяти HNSW (384d)", size=Pt(12), bold=True, color=DARK_BLUE)

    mem_table = slide.shapes.add_table(3, 4, Cm(1.5), Cm(4.2), Cm(11.5), Cm(3.5))
    tbl = mem_table.table
    for i, w in enumerate([Cm(3), Cm(2.8), Cm(2.8), Cm(2.9)]):
        tbl.columns[i].width = w
    mem_headers = ["Векторов", "Данные", "Граф", "Итого"]
    mem_rows = [
        ("10 000", "15.4 МБ", "27.1 МБ", "42.5 МБ"),
        ("50 000", "76.8 МБ", "139.4 МБ", "216.2 МБ"),
    ]
    for i, h in enumerate(mem_headers):
        cell = tbl.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
    for r, row_data in enumerate(mem_rows):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            if c == 3:
                p.font.bold = True
                p.font.color.rgb = BLUE

    add_accent_card(slide, Cm(1.5), Cm(8.2), Cm(11.5), Cm(2.5), GREEN)
    add_text_box(slide, Cm(2.3), Cm(8.5), Cm(10), Cm(2),
                 "Оптимизация: сегментированная архитектура\n"
                 "с memory-mapping (mmap) снижает нагрузку на RAM",
                 size=Pt(11), color=BLACK)

    # Functional tests section
    add_text_box(slide, Cm(14), Cm(3.2), Cm(10), Cm(0.8),
                 "Функциональное тестирование", size=Pt(12), bold=True, color=DARK_BLUE)

    add_metric_card(slide, Cm(14), Cm(4.2), Cm(4.8), Cm(2.5), "74+", "JVM-тестов", BLUE)
    add_metric_card(slide, Cm(19.3), Cm(4.2), Cm(4.8), Cm(2.5), "12/12", "сценариев ТЗ", GREEN)

    func_tests = [
        "Семантический поиск",
        "Фоновая индексация",
        "Управление альбомами",
        "Исключение каталогов",
        "Автономная работа (офлайн)",
        "Энергосбережение (<20%)",
    ]
    for i, item in enumerate(func_tests):
        y = Cm(7.5 + i * Cm(0.065))
        txBox = slide.shapes.add_textbox(Cm(14.5), Cm(7.5 + i * 1.5), Cm(9), Cm(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "\u2713  "
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = GREEN
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(11)
        run2.font.color.rgb = BLACK

    # =========== СЛАЙД 14: Заключение ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)
    add_title(slide, "Заключение")
    add_slide_number(slide, 14)

    conclusions = [
        ("Разработано полностью автономное кроссплатформенное приложение "
         "для семантического поиска без передачи данных на серверы", BLUE),
        ("Реализована двухуровневая ML-архитектура: CLIP для массовой "
         "индексации + VLM для глубокого анализа", PURPLE),
        ("Собственный HNSW-индекс с сегментированным хранением: "
         "p95 < 2 мс, recall@10 \u2265 0.95", GREEN),
        ("Все требования ТЗ выполнены, подтверждены 74+ тестами "
         "и бенчмарк-сюитой из 17 тестов", ORANGE),
    ]

    for i, (text, color) in enumerate(conclusions):
        y = Cm(3.5 + i * 3.2)
        add_accent_card(slide, Cm(1.5), y, Cm(22), Cm(2.7), color)
        add_text_box(slide, Cm(2.5), y + Cm(0.5), Cm(20.5), Cm(2),
                     text, size=Pt(13), color=BLACK)

    # =========== СЛАЙД 15: Спасибо за внимание ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_gradient_bar(slide)

    # Decorative circles
    for (cx, cy, sz, alpha) in [(Cm(2), Cm(13), Cm(6), 0xF0), (Cm(20), Cm(2), Cm(4), 0xF5)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(alpha, alpha, 0xFF)
        c.line.fill.background()

    add_text_box(slide, Cm(2), Cm(7), Cm(21), Cm(2.5),
                 "Спасибо за внимание",
                 size=Pt(30), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Cm(2), Cm(10.5), Cm(21), Cm(2),
                 "Миняев Михаил Игоревич  \u00b7  ИКБО-06-22",
                 size=Pt(13), color=GRAY, alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "ВКР_Миняев_МИ.pptx")
    prs.save(output_path)
    print(f"\u2713 Презентация сохранена: {output_path}")
    print(f"\u2713 Слайдов: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
