"""
Генератор PPTX-презентации ВКР — Recall (v4 - улучшенный слайд стека)
Запуск: source .venv/bin/activate && python docs/generate_pptx_v4.py
Результат: docs/ВКР_Миняев_МИ_v4.pptx
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE = RGBColor(0x00, 0x3D, 0xA5)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x7A)
ACCENT = RGBColor(0x00, 0x66, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x77)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

SLIDE_WIDTH = Cm(25.4)
SLIDE_HEIGHT = Cm(19.05)


def add_gradient_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.45)
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = DARK_BLUE
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = ACCENT
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Cm(2), SLIDE_HEIGHT - Cm(1.2), Cm(1.5), Cm(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Cm(2), Cm(1.1), Cm(21), Cm(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.5), Cm(3.5), Cm(0.12)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xFC)
    shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xEE)
    shape.line.width = Pt(0.75)
    return shape


def add_white_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xEE)
    shape.line.width = Pt(0.75)
    return shape


def add_text(slide, left, top, width, height, text, size=Pt(14),
             bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(4)
    return tf


def add_placeholder(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
    shape.line.dash_style = 2
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def add_metric(slide, left, top, width, value, desc, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Cm(2.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xEE)
    shape.line.width = Pt(0.75)
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Cm(0.15), top + Cm(0.08), width - Cm(0.3), Cm(0.13)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER


def add_tech_item(slide, x, y, letter, name, desc, bg_color):
    """Add a tech stack item with colored icon showing first letter."""
    # Icon circle with letter
    icon = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(1.6), Cm(1.6)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = bg_color
    icon.line.fill.background()
    tf = icon.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Name
    txBox = slide.shapes.add_textbox(x + Cm(2.0), y + Cm(0.0), Cm(8.5), Cm(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Description
    txBox2 = slide.shapes.add_textbox(x + Cm(2.0), y + Cm(0.8), Cm(8.5), Cm(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ===== СЛАЙД 1: Титульный =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_text(slide, Cm(2), Cm(2), Cm(21), Cm(2.2),
             "МИНОБРНАУКИ РОССИИ\n"
             "ФГБОУ ВО «МИРЭА — Российский технологический университет»\n"
             "Институт информационных технологий",
             size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(2), Cm(5.5), Cm(21), Cm(5),
             "Кроссплатформенное мобильное\nприложение для семантического\n"
             "анализа медиатеки с использованием\nвстраиваемых LLM",
             size=Pt(24), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(2), Cm(12.5), Cm(21), Cm(1),
             "студент группы ИКБО-06-22  Миняев М.И.",
             size=Pt(13), color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(2), Cm(14), Cm(21), Cm(1),
             "руководитель: [ФИО руководителя]",
             size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)

    # ===== СЛАЙД 2: Актуальность =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Актуальность")
    add_slide_number(slide, 2)

    facts = [
        ("4,88 млрд", "пользователей смартфонов (2024)", BLUE),
        ("952+", "фотографий на устройстве в среднем", ACCENT),
        ("70+ TOPS", "мощность мобильных NPU (2023–2025)", PURPLE),
        ("0", "офлайн-аналогов на рынке РФ", GREEN),
    ]
    for i, (val, label, color) in enumerate(facts):
        y = Cm(3.2 + i * 3.1)
        add_card(slide, Cm(1.5), y, Cm(11.5), Cm(2.7))
        add_text(slide, Cm(2.5), y + Cm(0.3), Cm(10), Cm(1.2),
                 val, size=Pt(20), bold=True, color=color)
        add_text(slide, Cm(2.5), y + Cm(1.5), Cm(10), Cm(0.9),
                 label, size=Pt(11), color=GRAY)

    right_cards = [
        ("Проблема", "Пользователь помнит содержимое\nфото, но ищет по дате и имени", RED_SOFT),
        ("Ограничение", "Облачные решения нарушают\nконфиденциальность данных", ORANGE),
        ("Научная новизна", "На рынке РФ нет приложений\nс локальным семантическим\nпоиском по медиатеке", GREEN),
    ]
    for i, (title, desc, color) in enumerate(right_cards):
        y = Cm(3.2 + i * 4.2)
        add_card(slide, Cm(14), y, Cm(10), Cm(3.7))
        add_text(slide, Cm(15), y + Cm(0.4), Cm(8.5), Cm(0.7),
                 title, size=Pt(11), bold=True, color=color)
        add_text(slide, Cm(15), y + Cm(1.3), Cm(8.5), Cm(2.2),
                 desc, size=Pt(12), color=BLACK)

    # ===== СЛАЙД 3: Цели и задачи =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Цели и задачи ВКР")
    add_slide_number(slide, 3)

    add_card(slide, Cm(1.5), Cm(3.2), Cm(11), Cm(9))
    add_text(slide, Cm(2.2), Cm(3.5), Cm(10), Cm(0.7),
             "ЦЕЛЬ", size=Pt(10), bold=True, color=BLUE)
    add_text(slide, Cm(2.2), Cm(4.3), Cm(10), Cm(7.5),
             "Разработать кроссплатформенное мобильное приложение "
             "для семантического поиска по локальной медиатеке "
             "с использованием встроенных LLM, обеспечивающее "
             "интеллектуальный поиск изображений при полной "
             "автономности и конфиденциальности данных",
             size=Pt(13), color=BLACK)

    add_text(slide, Cm(13.5), Cm(3.2), Cm(5), Cm(0.7),
             "ЗАДАЧИ", size=Pt(10), bold=True, color=BLUE)
    tasks = [
        "Анализ предметной области",
        "Обоснование необходимости автоматизации",
        "Формулировка требований к системе",
        "Разработка архитектуры и информ. модели",
        "Реализация программных компонентов",
        "Тестирование программного решения",
    ]
    for i, task in enumerate(tasks):
        y = Cm(4.2 + i * 2.0)
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(13.5), y, Cm(0.9), Cm(0.9))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        add_text(slide, Cm(14.8), y + Cm(0.05), Cm(9.5), Cm(0.9),
                 task, size=Pt(12), color=BLACK)

    # ===== СЛАЙД 4: Системы-аналоги =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Системы-аналоги")
    add_slide_number(slide, 4)

    headers = ["Система", "Семантический\nпоиск", "Полный\nофлайн", "Кросс-\nплатформ."]
    rows_data = [
        ("Google Photos", "\u2713", "\u2717", "\u2713"),
        ("Apple Photos", "\u25D0", "\u2717", "\u2717"),
        ("Simple Gallery Pro", "\u2717", "\u2713", "\u2717"),
        ("Immich", "\u2713", "\u2717", "\u2713"),
        ("Ente Photos", "\u2717", "\u2717", "\u2713"),
        ("Recall (наша)", "\u2713", "\u2713", "\u2713"),
    ]
    table_shape = slide.shapes.add_table(
        len(rows_data) + 1, 4, Cm(1.5), Cm(3.5), Cm(22), Cm(10.5)
    )
    table = table_shape.table
    for i, w in enumerate([Cm(7.5), Cm(5), Cm(4.5), Cm(5)]):
        table.columns[i].width = w
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if val == "\u2713":
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif val == "\u2717":
                p.font.color.rgb = RED_SOFT
            if r == len(rows_data) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
                p.font.bold = True
    add_text(slide, Cm(1.5), Cm(14.8), Cm(22), Cm(1),
             "Ни одно решение не объединяет семантический поиск + офлайн + кроссплатформенность",
             size=Pt(12), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # ===== СЛАЙД 5: Требования =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Требования к системе")
    add_slide_number(slide, 5)
    add_placeholder(slide, Cm(1.5), Cm(3.2), Cm(22), Cm(13.5),
                    "[Вставить: Диаграмма вариантов использования (Use Case)]\n\n"
                    "Источник: рисунок 2.5 из проектной практики\n\n"
                    "Акторы: Пользователь, Система")

    # ===== СЛАЙД 6: Архитектура =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Архитектура системы")
    add_slide_number(slide, 6)
    add_placeholder(slide, Cm(1.5), Cm(3.2), Cm(15), Cm(13.5),
                    "[Вставить: Компонентная диаграмма C4]\n\n"
                    "Источник: рисунок 2.3 из проектной практики\n\n"
                    "UI \u2192 Feature \u2192 Core")
    principles = ["Clean Architecture", "Модульная декомпозиция",
                  "0 серверных компонентов", "MVI (однонаправленный поток)", "Offline-first"]
    for i, text in enumerate(principles):
        y = Cm(3.5 + i * 2.5)
        add_card(slide, Cm(17), y, Cm(7.2), Cm(2.1))
        add_text(slide, Cm(17.8), y + Cm(0.5), Cm(6.2), Cm(1.1),
                 text, size=Pt(11), bold=True, color=DARK_BLUE)

    # ===== СЛАЙД 7: Информационная модель =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Информационная модель (схема БД)")
    add_slide_number(slide, 7)
    add_placeholder(slide, Cm(1.5), Cm(3.2), Cm(14.5), Cm(13.5),
                    "[Вставить: ER-диаграмма]\n\n"
                    "Источник: рисунок 2.9 из проектной практики\n\n"
                    "6 сущностей: MediaFile, EmbeddingBatch,\n"
                    "Album, AlbumEntry, ExcludedPath, SearchQuery")
    decisions = [
        ("Хранение векторов", "В бинарных батч-файлах,\nне в SQLite"),
        ("Доступ к вектору", "O(1): offset = k \u00d7 d \u00d7 4 bytes"),
        ("СУБД", "SQLite через Room 2.8"),
    ]
    for i, (title, desc) in enumerate(decisions):
        y = Cm(3.5 + i * 4.2)
        add_card(slide, Cm(17), y, Cm(7.2), Cm(3.7))
        add_text(slide, Cm(17.8), y + Cm(0.4), Cm(6), Cm(0.7),
                 title, size=Pt(10), bold=True, color=ACCENT)
        add_text(slide, Cm(17.8), y + Cm(1.3), Cm(6), Cm(2),
                 desc, size=Pt(12), color=BLACK)

    # ===== СЛАЙД 8: Математическое обеспечение =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Математическое обеспечение")
    add_slide_number(slide, 8)

    # Section 1
    y1 = Cm(3.2)
    add_card(slide, Cm(1.5), y1, Cm(22), Cm(3.6))
    add_text(slide, Cm(2.3), y1 + Cm(0.3), Cm(20), Cm(0.7),
             "1. Генерация дескрипторов (CLIP)", size=Pt(13), bold=True, color=BLUE)
    f1 = add_white_card(slide, Cm(2.3), y1 + Cm(1.2), Cm(20.5), Cm(1.9))
    tf = f1.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for txt, sz, clr, bld, ital in [
        ("f", Pt(15), BLACK, False, True), ("img", Pt(9), GRAY, False, False),
        (" : Image \u2192 \u211d", Pt(14), BLACK, False, False), ("d", Pt(9), BLACK, False, False),
        ("              ", Pt(14), BLACK, False, False),
        ("f", Pt(15), BLACK, False, True), ("text", Pt(9), GRAY, False, False),
        (" : Text \u2192 \u211d", Pt(14), BLACK, False, False), ("d", Pt(9), BLACK, False, False),
        ("              ", Pt(14), BLACK, False, False),
        ("d = 384 / 512", Pt(13), ACCENT, True, False),
    ]:
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.italic = ital

    # Section 2
    y2 = Cm(7.3)
    add_card(slide, Cm(1.5), y2, Cm(22), Cm(3.6))
    add_text(slide, Cm(2.3), y2 + Cm(0.3), Cm(20), Cm(0.7),
             "2. Поиск ближайшего соседа", size=Pt(13), bold=True, color=PURPLE)
    f2 = add_white_card(slide, Cm(2.3), y2 + Cm(1.2), Cm(20.5), Cm(1.9))
    tf = f2.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for txt, sz, clr, bld in [
        ("sim(q\u0302, v\u0302", Pt(15), BLACK, False), ("i", Pt(9), BLACK, False),
        (") = q\u0302 \u00b7 v\u0302", Pt(15), BLACK, False), ("i", Pt(9), BLACK, False),
        ("              ", Pt(14), BLACK, False),
        ("Results = { i  |  sim \u2265 \u03c4 }", Pt(14), BLACK, False),
    ]:
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.color.rgb = clr
        run.font.bold = bld
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "Косинусное сходство нормированных векторов = скалярное произведение"
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY

    # Section 3
    y3 = Cm(11.4)
    add_card(slide, Cm(1.5), y3, Cm(22), Cm(4.8))
    add_text(slide, Cm(2.3), y3 + Cm(0.3), Cm(20), Cm(0.7),
             "3. HNSW (Hierarchical Navigable Small World)", size=Pt(13), bold=True, color=GREEN)
    f3l = add_white_card(slide, Cm(2.3), y3 + Cm(1.2), Cm(10), Cm(3.0))
    tf = f3l.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Линейный перебор: "; r.font.size = Pt(13); r.font.color.rgb = GRAY
    r = p.add_run(); r.text = "O(n)"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RED_SOFT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = "\u2193"; r.font.size = Pt(14); r.font.color.rgb = LIGHT_GRAY
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run(); r.text = "HNSW: "; r.font.size = Pt(13); r.font.color.rgb = GRAY
    r = p3.add_run(); r.text = "O(log n)"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = GREEN

    f3r = add_white_card(slide, Cm(13), y3 + Cm(1.2), Cm(10), Cm(3.0))
    tf = f3r.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.6)
    p = tf.paragraphs[0]
    p.text = "Многоуровневый навигируемый граф"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    p2.text = "M = 16     efConstruction = 200     efSearch = 50"
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT

    # ===== СЛАЙД 9: Инструменты (НОВЫЙ ДИЗАЙН) =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Инструментальные средства разработки")
    add_slide_number(slide, 9)

    # Tools with colored icon squares
    tools = [
        # (letter, name, description, icon_color)
        ("K", "Kotlin 2.2", "Основной язык — типобезопасность, корутины, null-safety", RGBColor(0x7F, 0x52, 0xFF)),
        ("KMP", "Kotlin Multiplatform", "Единая бизнес-логика для Android и iOS без мостов", RGBColor(0x00, 0x95, 0xD5)),
        ("JC", "Jetpack Compose", "Декларативный UI-фреймворк, Material 3, анимации", RGBColor(0x41, 0x85, 0xF4)),
        ("R", "Room 2.8", "Типобезопасная ORM поверх SQLite с миграциями", RGBColor(0xA4, 0xC6, 0x39)),
        ("TF", "TensorFlow Lite 2.17", "Инференс CLIP-модели на NPU/GPU через NNAPI", RGBColor(0xFF, 0x6F, 0x00)),
        ("H", "HNSW (Kotlin)", "Собственный индекс: векторный поиск за O(log n)", RGBColor(0x10, 0xB9, 0x81)),
        ("W", "WorkManager", "Фоновая индексация с учётом заряда и температуры", RGBColor(0x00, 0x89, 0x7B)),
        ("DI", "Hilt (DI)", "Внедрение зависимостей — модульность и тестируемость", RGBColor(0xE9, 0x1E, 0x63)),
        ("M", "MVI", "Однонаправленный поток — предотвращает race conditions", RGBColor(0x9C, 0x27, 0xB0)),
        ("T", "kotlin.test + JUnit", "Единые тесты для обеих платформ, 74+ тестов", RGBColor(0x79, 0x55, 0x48)),
        ("AS", "Android Studio", "IDE с Memory Profiler и Energy Profiler", RGBColor(0x3D, 0xDC, 0x84)),
        ("G", "Git", "Контроль версий и история изменений", RGBColor(0xF0, 0x50, 0x33)),
    ]

    for i, (letter, name, desc, icon_color) in enumerate(tools):
        col = i % 2
        row = i // 2
        x = Cm(1.3) + col * Cm(12.3)
        y = Cm(3.1) + row * Cm(2.35)
        add_tech_item(slide, x, y, letter, name, desc, icon_color)

    # ===== СЛАЙД 10: Программное решение 1 =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Программное решение")
    add_slide_number(slide, 10)
    scr1 = [
        ("[Экран приветствия]\n\nрис. 4.1 тех. практики", "Рис. 1 \u2014 Onboarding"),
        ("[Главный экран]\n\nрис. 4.2 тех. практики", "Рис. 2 \u2014 Timeline"),
        ("[Альбомы]\n\nfeature/albums", "Рис. 3 \u2014 Альбомы"),
    ]
    for i, (txt, cap) in enumerate(scr1):
        x = Cm(1.2) + i * Cm(8)
        add_placeholder(slide, x, Cm(3.2), Cm(7.5), Cm(11.5), txt)
        add_text(slide, x, Cm(15), Cm(7.5), Cm(0.8),
                 cap, size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

    # ===== СЛАЙД 11: Программное решение 2 =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Программное решение")
    add_slide_number(slide, 11)
    scr2 = [
        ("[Поисковый запрос]\n\nрис. 4.4 тех. практики", "Рис. 4 \u2014 Поиск"),
        ("[Результаты]\n\nрис. 4.5 тех. практики", "Рис. 5 \u2014 Результаты"),
        ("[Исключения]\n\nрис. 4.3 тех. практики", "Рис. 6 \u2014 Исключения"),
    ]
    for i, (txt, cap) in enumerate(scr2):
        x = Cm(1.2) + i * Cm(8)
        add_placeholder(slide, x, Cm(3.2), Cm(7.5), Cm(11.5), txt)
        add_text(slide, x, Cm(15), Cm(7.5), Cm(0.8),
                 cap, size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

    # ===== СЛАЙД 12: Бенчмарки =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Тестирование \u2014 бенчмарки производительности")
    add_slide_number(slide, 12)
    add_placeholder(slide, Cm(1.5), Cm(3.2), Cm(10.8), Cm(6),
                    "[График: Search Latency p95]\nHNSW vs LinearScan (1K\u201350K)\n\n"
                    "docs/charts-real.html \u2014 Chart 1")
    add_placeholder(slide, Cm(12.8), Cm(3.2), Cm(10.8), Cm(6),
                    "[График: Recall@10 vs Index Size]\nM=16, efSearch=50\n\n"
                    "docs/charts-real.html \u2014 Chart 2")
    metrics = [
        ("< 2 мс", "p95 @ 50K", BLUE),
        ("\u2265 0.95", "Recall@10", PURPLE),
        ("274 v/s", "throughput", GREEN),
        ("16 мс", "serialize 10K", ORANGE),
    ]
    for i, (v, d, c) in enumerate(metrics):
        add_metric(slide, Cm(1.5) + i * Cm(5.8), Cm(9.8), Cm(5.3), v, d, c)
    add_text(slide, Cm(1.5), Cm(13.3), Cm(22), Cm(2.5),
             "17 бенчмарков  \u00b7  384d  \u00b7  seed=42  \u00b7  100 queries\n"
             "Требование ТЗ: < 50 мс @ 50K \u2014 выполнено с запасом 25\u00d7",
             size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)

    # ===== СЛАЙД 13: Память + функц. тесты =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Тестирование \u2014 память и функциональные тесты")
    add_slide_number(slide, 13)
    add_text(slide, Cm(1.5), Cm(3.2), Cm(11), Cm(0.7),
             "Потребление памяти HNSW (384d)", size=Pt(12), bold=True, color=DARK_BLUE)
    mem_tbl = slide.shapes.add_table(3, 4, Cm(1.5), Cm(4.2), Cm(11.5), Cm(3.2))
    tbl = mem_tbl.table
    for i, w in enumerate([Cm(3), Cm(2.8), Cm(2.8), Cm(2.9)]):
        tbl.columns[i].width = w
    for i, h in enumerate(["Векторов", "Данные", "Граф", "Итого"]):
        cell = tbl.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
    for r, row in enumerate([("10 000", "15.4 МБ", "27.1 МБ", "42.5 МБ"),
                              ("50 000", "76.8 МБ", "139.4 МБ", "216.2 МБ")]):
        for c, v in enumerate(row):
            cell = tbl.cell(r+1, c)
            cell.text = v
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11); p.alignment = PP_ALIGN.CENTER
            if c == 3: p.font.bold = True; p.font.color.rgb = BLUE

    add_card(slide, Cm(1.5), Cm(7.8), Cm(11.5), Cm(2.2))
    add_text(slide, Cm(2.3), Cm(8.1), Cm(10.5), Cm(1.7),
             "Оптимизация: сегментированная архитектура\n"
             "с memory-mapping (mmap) снижает нагрузку на RAM",
             size=Pt(11), color=BLACK)

    add_text(slide, Cm(14), Cm(3.2), Cm(10), Cm(0.7),
             "Функциональное тестирование", size=Pt(12), bold=True, color=DARK_BLUE)
    add_metric(slide, Cm(14), Cm(4.2), Cm(4.8), "74+", "JVM-тестов", BLUE)
    add_metric(slide, Cm(19.3), Cm(4.2), Cm(4.8), "12/12", "сценариев ТЗ", GREEN)
    func_items = ["Семантический поиск", "Фоновая индексация",
                  "Управление альбомами", "Исключение каталогов",
                  "Автономная работа (офлайн)", "Энергосбережение (<20%)"]
    for i, item in enumerate(func_items):
        txBox = slide.shapes.add_textbox(Cm(14.5), Cm(7.5 + i * 1.4), Cm(9), Cm(1.1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "\u2713  "; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREEN
        r = p.add_run(); r.text = item; r.font.size = Pt(11); r.font.color.rgb = BLACK

    # ===== СЛАЙД 14: Заключение =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_title(slide, "Заключение")
    add_slide_number(slide, 14)
    conclusions = [
        "Разработано полностью автономное кроссплатформенное приложение "
        "для семантического поиска без передачи данных на серверы",
        "Реализована двухуровневая ML-архитектура: CLIP для массовой "
        "индексации + Vision-Language модель для глубокого анализа",
        "Собственный HNSW-индекс с сегментированным хранением: "
        "p95 < 2 мс, recall@10 \u2265 0.95",
        "Все требования ТЗ выполнены, подтверждены 74+ тестами "
        "и бенчмарк-сюитой из 17 тестов",
    ]
    colors = [BLUE, PURPLE, GREEN, ORANGE]
    for i, (text, color) in enumerate(zip(conclusions, colors)):
        y = Cm(3.4 + i * 3.3)
        add_card(slide, Cm(1.5), y, Cm(22), Cm(2.8))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.3), y + Cm(1.0), Cm(0.4), Cm(0.4))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(slide, Cm(3.2), y + Cm(0.5), Cm(19.5), Cm(2),
                 text, size=Pt(13), color=BLACK)

    # ===== СЛАЙД 15: Спасибо =====
    slide = prs.slides.add_slide(blank)
    add_gradient_bar(slide)
    add_text(slide, Cm(2), Cm(7), Cm(21), Cm(2),
             "Спасибо за внимание",
             size=Pt(30), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(2), Cm(10.5), Cm(21), Cm(1.5),
             "Миняев Михаил Игоревич  \u00b7  ИКБО-06-22",
             size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

    out = os.path.join(os.path.dirname(__file__), "ВКР_Миняев_МИ_v4.pptx")
    prs.save(out)
    print(f"\u2713 Сохранено: {out}")
    print(f"\u2713 Слайдов: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
